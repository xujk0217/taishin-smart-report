"""Tests for agent-generated Python PPTX rendering."""

from __future__ import annotations

import json
from typing import Any, AsyncGenerator, AsyncIterable

import pytest
from pptx import Presentation
from openpyxl import Workbook
from strands.models import Model

from lobster_runtime.artifact_validator import validate_artifacts
from lobster_runtime.presentation_agent import AgentPresentationRuntime
from lobster_runtime.presentation_contracts import EvidencePacketV2, PresentationBlueprint, RenderArtifactManifest
from lobster_runtime.render_code_guard import validate_renderer_source
from lobster_runtime.smart_report_pptx import Deck, RenderingContext
from lobster_runtime.universal_presentation_pipeline import UniversalPresentationPipeline
from lobster_runtime.universal_pipeline_contracts import (
    AnalysisFeasibilityPlan,
    BlueprintStageOutput,
    DataIntelligenceReport,
    VerifiedAnalysisNarrative,
)


RENDERER_SOURCE = '''
from lobster_runtime.smart_report_pptx import Deck

def render(ctx):
    deck = Deck.from_context(ctx)
    slide = deck.add_slide(role="cover")
    slide.add_title("營運表現摘要", box=(0.7, 0.7, 10.5, 0.8))
    slide.add_claim(claim_ref="claim_summary", box=(0.9, 1.8, 7.5, 1.2))

    slide = deck.add_slide(role="content")
    slide.add_title("區域比較", box=(0.7, 0.4, 9.0, 0.6))
    slide.add_metric_card(metric_ref="metric_revenue", label="營收", box=(0.7, 1.2, 2.5, 1.0))
    slide.add_chart(chart_ref="chart_region", box=(0.7, 2.5, 6.0, 3.5))
    slide.add_table(table_ref="table_region", box=(7.2, 2.5, 4.8, 2.4))

    slide = deck.add_slide(role="back-cover")
    slide.add_title("後續行動", box=(0.7, 0.7, 10.5, 0.8))
    slide.add_free_text("請依據已驗證結果安排行動。", box=(0.9, 1.8, 6.5, 1.0))
    return deck.save()
'''.strip()

PIPELINE_RENDERER_SOURCE = '''
from lobster_runtime.smart_report_pptx import Deck

def render(ctx):
    deck = Deck.from_context(ctx)
    slide = deck.add_slide(role="cover")
    slide.add_title("營運表現摘要", box=(0.7, 0.7, 10.5, 0.8))
    slide.add_claim(claim_ref="claim_1_summary", box=(0.9, 1.8, 7.5, 1.2))
    slide.add_metric_card(metric_ref="metric_1_營收", label="營收", box=(0.9, 3.0, 2.8, 1.0))
    slide.add_chart(chart_ref="chart_1_營收", box=(4.0, 2.7, 6.2, 3.2))
    slide.add_table(table_ref="table_1_營收", box=(0.9, 4.4, 3.0, 1.8))
    return deck.save()
'''.strip()

NO_RETURN_RENDERER_SOURCE = '''
from lobster_runtime.smart_report_pptx import Deck

def render(ctx):
    deck = Deck.from_context(ctx)
    slide = deck.add_slide(role="cover")
    slide.add_title("營運表現摘要", box=(0.7, 0.7, 10.5, 0.8))
    slide.add_claim(claim_ref="claim_summary", box=(0.9, 1.8, 7.5, 1.2))
'''.strip()

NATIVE_PPTX_RENDERER_SOURCE = '''
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from openpyxl import Workbook

def render(ctx):
    prs = ctx.new_presentation()
    wb = ctx.new_workbook()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.5), Inches(0.7))
    title_box.text_frame.text = "營運表現摘要"
    claim = ctx.claims["claim_summary"]
    metric = ctx.metrics["metric_revenue"]
    claim_text = claim.text.replace("{{metric_revenue}}", f"{metric.value:,.0f}{metric.unit}")
    claim_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.0), Inches(0.9))
    claim_box.text_frame.text = claim_text

    chart = ctx.charts["chart_region"]
    chart_data = CategoryChartData()
    chart_data.categories = chart.categories
    for series in chart.series:
        chart_data.add_series(series.name, series.values)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(2.8), Inches(6.4), Inches(3.2), chart_data)

    sheet = wb.create_sheet("chart_region")
    sheet.append(["category", "營收"])
    for category, value in zip(chart.categories, chart.series[0].values):
        sheet.append([category, value])
    return ctx.save_artifacts(prs, wb, chart_count=1, table_count=0, evidence_refs_used=["claim_summary", "metric_revenue", "chart_region"])
'''.strip()

DICT_STYLE_NATIVE_RENDERER_SOURCE = '''
from pptx.util import Inches

def render(ctx):
    prs = ctx.new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    metric = ctx.metrics["metric_revenue"]
    claim = ctx.claims["claim_summary"]
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(8.0), Inches(0.8))
    box.text_frame.text = "營運表現摘要"
    claim_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.0), Inches(0.8))
    claim_box.text_frame.text = claim["text"].replace("{{metric_revenue}}", f"{metric['value']:,.0f}{metric.get('unit', '')}")
    return ctx.save_artifacts(prs, evidence_refs_used=["claim_summary", "metric_revenue"])
'''.strip()


class ScriptedRendererModel(Model):
    def __init__(self, source: str = RENDERER_SOURCE) -> None:
        self.source = source
        self.config: dict[str, Any] = {"model_id": "scripted-renderer-test"}

    def get_config(self) -> dict[str, Any]:
        return dict(self.config)

    def update_config(self, **model_config: Any) -> None:
        self.config.update(model_config)

    async def stream(self, messages: list[dict[str, Any]], tool_specs: list[dict[str, Any]] | None = None, system_prompt: str | None = None, **_: Any) -> AsyncIterable[dict[str, Any]]:
        names = {spec["name"] for spec in tool_specs or []}
        if "DataIntelligenceReport" in names:
            output = {
                "stage": "data-intelligence",
                "status": "passed",
                "data_structure_summary": "Workbook contains a region table with numeric revenue.",
                "semantic_notes": ["Region is categorical and revenue is numeric."],
                "relationship_hypotheses": [],
                "quality_findings": [],
                "usable_data_notes": ["The table can support a summary presentation."],
            }
        elif "AnalysisFeasibilityPlan" in names:
            output = {
                "stage": "analysis-feasibility",
                "status": "passed",
                "accepted_analyses": [{
                    "analysis_id": "analysis_region_revenue",
                    "question": "How does revenue compare by region?",
                    "method": "Sum numeric revenue columns and compare categories.",
                    "required_columns": ["區域", "營收"],
                    "feasible": True,
                    "rationale": "The workbook has visible category and numeric values.",
                    "recommended_visual": "column chart",
                }],
                "rejected_analyses": [],
                "known_limits": [],
                "questions_for_user": [],
            }
        elif "VerifiedAnalysisNarrative" in names:
            output = {
                "stage": "verified-analysis",
                "status": "passed",
                "insight_summaries": ["Use the generated evidence packet for all numeric content."],
                "caveats": [],
                "evidence_usage_notes": ["Claims and charts must reference evidence IDs."],
            }
        elif "BlueprintStageOutput" in names:
            output = {
                "stage": "presentation-design",
                "status": "passed",
                "blueprint": {
                    "blueprint_version": "presentation-blueprint-v1",
                    "title": "營運表現摘要",
                    "slides": [{
                        "slide_id": "s1",
                        "role": "cover",
                        "intent": "Summarize evidence-backed result.",
                        "layout_strategy": "claim led title slide",
                        "elements": [
                            {"element_id": "e1", "type": "title", "box": {"x": 0.7, "y": 0.7, "w": 10.5, "h": 0.8}, "text": "營運表現摘要"},
                            {"element_id": "e2", "type": "claim", "box": {"x": 0.9, "y": 1.8, "w": 7.5, "h": 1.2}, "claim_ref": "claim_1_summary"},
                            {"element_id": "e3", "type": "chart", "box": {"x": 0.9, "y": 3.0, "w": 6.0, "h": 3.0}, "chart_ref": "chart_1_營收"},
                        ],
                    }],
                },
                "design_notes": ["Use template if present."],
            }
        elif "PythonRendererProgram" in names:
            output = {
            "program_version": "python-renderer-program-v1",
            "entrypoint": "render",
            "source_code": self.source,
            "rationale": "Use controlled SDK calls and evidence refs only.",
            }
        else:
            raise AssertionError(f"unexpected tools: {sorted(names)}")
        encoded = json.dumps(output, ensure_ascii=False)
        yield {"messageStart": {"role": "assistant"}}
        yield {"contentBlockStart": {"start": {"toolUse": {"toolUseId": "renderer-output-1", "name": "PythonRendererProgram"}}}}
        yield {"contentBlockDelta": {"delta": {"toolUse": {"input": encoded}}}}
        yield {"contentBlockStop": {}}
        yield {"messageStop": {"stopReason": "tool_use"}}
        yield {"metadata": {"usage": {"inputTokens": 1, "outputTokens": 1, "totalTokens": 2}, "metrics": {"latencyMs": 1}}}

    async def structured_output(self, *_: Any, **__: Any) -> AsyncGenerator[dict[str, Any], None]:
        raise AssertionError("runtime must use the current structured_output_model invocation API")
        if False:
            yield {}


def evidence_payload() -> dict[str, Any]:
    return {
        "packet_id": "evp-demo-v2",
        "metrics": [{
            "metric_id": "metric_revenue",
            "label": "營收",
            "value": 1250000,
            "unit": "元",
            "display_format": "currency",
            "source_refs": ["source_sheet_region"],
            "calculation": "sum revenue by selected period",
        }],
        "charts": [{
            "chart_id": "chart_region",
            "title": "區域營收",
            "chart_type": "column",
            "categories": ["北", "中", "南"],
            "series": [{"name": "營收", "values": [600000, 350000, 300000]}],
            "metric_refs": ["metric_revenue"],
        }],
        "claims": [{
            "claim_id": "claim_summary",
            "text": "本期營收為 {{metric_revenue}}，後續應聚焦主要區域。",
            "metric_refs": ["metric_revenue"],
            "chart_refs": ["chart_region"],
            "source_refs": ["source_sheet_region"],
        }],
        "tables": [{
            "table_id": "table_region",
            "title": "區域營收表",
            "headers": ["區域", "營收"],
            "rows": [["北", 600000], ["中", 350000], ["南", 300000]],
            "source_refs": ["source_sheet_region"],
        }],
    }


def blueprint_payload() -> dict[str, Any]:
    return {
        "blueprint_version": "presentation-blueprint-v1",
        "title": "營運表現摘要",
        "slides": [{
            "slide_id": "s1",
            "role": "cover",
            "intent": "建立報告主題",
            "layout_strategy": "large title with claim",
            "elements": [{"element_id": "e1", "type": "claim", "box": {"x": 0.9, "y": 1.8, "w": 7.5, "h": 1.2}, "claim_ref": "claim_summary"}],
        }],
    }


def test_agent_generated_renderer_creates_pptx_and_xlsx(tmp_path) -> None:
    runtime = AgentPresentationRuntime(ScriptedRendererModel())
    result = runtime.generate(
        prompt="產生一份營運摘要簡報",
        blueprint=PresentationBlueprint.model_validate(blueprint_payload()),
        evidence=EvidencePacketV2.model_validate(evidence_payload()),
        output_dir=tmp_path,
    )

    assert result.validation_report.status == "final"
    assert result.manifest.slide_count == 3
    assert result.manifest.chart_count == 1
    assert result.manifest.table_count == 1
    assert result.template_profile.source == "default"
    assert result.manifest.evidence_refs_used == ["chart_region", "claim_summary", "metric_revenue", "table_region"]
    assert Presentation(result.manifest.pptx_path).slides


def test_agent_renderer_applies_uploaded_template(tmp_path) -> None:
    template = tmp_path / "template.pptx"
    template_deck = Presentation()
    template_deck.slides.add_slide(template_deck.slide_layouts[0])
    template_deck.slides.add_slide(template_deck.slide_layouts[1])
    template_deck.save(template)
    runtime = AgentPresentationRuntime(ScriptedRendererModel())
    result = runtime.generate(
        prompt="套用上傳模板產生簡報",
        blueprint=blueprint_payload(),
        evidence=evidence_payload(),
        output_dir=tmp_path,
        template_path=template,
        file_stem="with-template",
    )

    assert result.template_profile.source == "uploaded"
    assert result.template_profile.template_path == str(template)
    assert result.validation_report.status == "final"
    assert result.manifest.slide_count == 3
    assert len(Presentation(result.manifest.pptx_path).slides) == 3


def test_agent_renderer_auto_saves_when_render_returns_none(tmp_path) -> None:
    runtime = AgentPresentationRuntime(ScriptedRendererModel(source=NO_RETURN_RENDERER_SOURCE))
    result = runtime.generate(
        prompt="產生一份營運摘要簡報",
        blueprint=blueprint_payload(),
        evidence=evidence_payload(),
        output_dir=tmp_path,
    )

    assert result.validation_report.status == "final"
    assert result.manifest.slide_count == 1


def test_agent_renderer_can_use_native_python_pptx(tmp_path) -> None:
    runtime = AgentPresentationRuntime(ScriptedRendererModel(source=NATIVE_PPTX_RENDERER_SOURCE))
    result = runtime.generate(
        prompt="產生一份營運摘要簡報",
        blueprint=blueprint_payload(),
        evidence=evidence_payload(),
        output_dir=tmp_path,
    )

    assert result.validation_report.status == "final"
    assert result.manifest.slide_count == 1
    assert result.manifest.chart_count == 1
    assert set(result.manifest.evidence_refs_used) == {"chart_region", "claim_summary", "metric_revenue"}


def test_agent_renderer_can_use_dict_style_evidence_access(tmp_path) -> None:
    runtime = AgentPresentationRuntime(ScriptedRendererModel(source=DICT_STYLE_NATIVE_RENDERER_SOURCE))
    result = runtime.generate(
        prompt="產生一份營運摘要簡報",
        blueprint=blueprint_payload(),
        evidence=evidence_payload(),
        output_dir=tmp_path,
    )

    assert result.validation_report.status == "final"
    assert result.manifest.slide_count == 1
    assert set(result.manifest.evidence_refs_used) == {"claim_summary", "metric_revenue"}


def test_rendering_context_evidence_indexes_are_dict_and_list_like(tmp_path) -> None:
    ctx = RenderingContext(
        evidence=EvidencePacketV2.model_validate(evidence_payload()),
        output_dir=tmp_path,
        file_stem="index",
    )

    assert ctx.metrics["metric_revenue"]["value"] == 1250000
    assert ctx.metrics[0]["id"] == "metric_revenue"
    assert [metric["id"] for metric in ctx.metrics] == ["metric_revenue"]
    assert ctx.metric("metric_revenue").get("name") == "營收"
    assert ctx.pptx_path.endswith("index.pptx")
    assert ctx.xlsx_path.endswith("index.xlsx")


def test_deck_tolerates_agent_positional_refs_and_bad_layout_index(tmp_path) -> None:
    ctx = RenderingContext(
        evidence=EvidencePacketV2.model_validate(evidence_payload()),
        output_dir=tmp_path,
        file_stem="positional",
    )
    deck = Deck.from_context(ctx)
    slide = deck.add_slide("content", layout_index=999)
    slide.add_title("備援標題", text="營運摘要", box=(0.5, 0.5, 5.0, 0.6))
    slide.add_metric_card("metric_revenue", box=(0.5, 1.2, 2.5, 1.0))
    slide.add_claim("claim_summary", box=(0.5, 2.4, 6.0, 1.0))
    slide.add_chart("chart_region", box=(0.5, 3.4, 6.0, 2.5))
    slide.add_table("table_region", box=(7.0, 3.4, 4.0, 2.0))

    manifest = deck.save()

    assert manifest.slide_count == 1
    assert manifest.chart_count == 1
    assert manifest.table_count == 1


def test_deck_tolerates_invalid_agent_box_values(tmp_path) -> None:
    ctx = RenderingContext(
        evidence=EvidencePacketV2.model_validate(evidence_payload()),
        output_dir=tmp_path,
        file_stem="bad-box",
    )
    deck = Deck.from_context(ctx)
    slide = deck.add_slide("content")
    slide.add_title("營運摘要", box="xxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxx")
    slide.add_text("後續行動", box=(0.5, 1.2, 3.0, 0.8))
    slide.add_metric_card("metric_revenue", box={"x": "bad", "y": 1, "w": 2, "h": 1})

    manifest = deck.save()

    assert manifest.slide_count == 1
    assert Presentation(manifest.pptx_path).slides


def test_context_can_preserve_template_sample_slides_for_editing(tmp_path) -> None:
    template = tmp_path / "template-components.pptx"
    template_deck = Presentation()
    slide = template_deck.slides.add_slide(template_deck.slide_layouts[0])
    slide.shapes.add_textbox(0, 0, 1000000, 500000).text_frame.text = "Template title"
    template_deck.save(template)
    ctx = RenderingContext(
        evidence=EvidencePacketV2.model_validate(evidence_payload()),
        output_dir=tmp_path,
        template_path=template,
        file_stem="preserve-template",
    )

    prs = ctx.new_presentation(preserve_template_slides=True)

    assert len(prs.slides) == 1
    assert any(
        getattr(shape, "has_text_frame", False) and "Template title" in shape.text_frame.text
        for shape in prs.slides[0].shapes
    )


def test_deck_tolerates_common_agent_aliases_and_ref_keywords(tmp_path) -> None:
    ctx = RenderingContext(
        evidence=EvidencePacketV2.model_validate(evidence_payload()),
        output_dir=tmp_path,
        file_stem="aliases",
    )
    deck = Deck.from_context(ctx)
    slide = deck.add_slide("content")
    slide.add_subtitle("重點摘要", box=(0.5, 0.5, 4.0, 0.6), color="#333333")
    slide.add_summary(text="請依據已驗證證據安排後續行動。", box=(0.5, 1.2, 5.0, 0.8))
    slide.add_bullet_list(["依據證據檢視重點", "追蹤資料品質"], box=(0.5, 2.0, 5.0, 1.0))
    slide.add_metric_card(id="metric_revenue", box=(0.5, 3.2, 2.5, 1.0), style="compact")
    slide.add_claim(ref="claim_summary", box=(3.3, 3.2, 4.0, 1.0))
    slide.add_chart(id="chart_region", box=(0.5, 4.4, 5.5, 2.2))
    slide.add_table(ref="table_region", box=(6.4, 4.4, 4.0, 2.0))

    manifest = deck.save()

    assert manifest.chart_count == 1
    assert manifest.table_count == 1
    assert set(manifest.evidence_refs_used) == {"chart_region", "claim_summary", "metric_revenue", "table_region"}


def test_artifact_validator_warns_but_does_not_block_rendered_evidence_numbers(tmp_path) -> None:
    ctx = RenderingContext(
        evidence=EvidencePacketV2.model_validate(evidence_payload()),
        output_dir=tmp_path,
        file_stem="numbers",
    )
    deck = Deck.from_context(ctx)
    slide = deck.add_slide("content")
    slide.add_claim("claim_summary", box=(0.5, 0.5, 6.0, 1.0))
    manifest = deck.save()

    report = validate_artifacts(RenderArtifactManifest.model_validate(manifest))

    assert report.status == "final"
    assert any(f.code == "POSSIBLE_UNTRACED_NUMBER" and f.severity == "warning" for f in report.findings)


def test_artifact_validator_blocks_chart_refs_without_native_chart(tmp_path) -> None:
    ctx = RenderingContext(
        evidence=EvidencePacketV2.model_validate(evidence_payload()),
        output_dir=tmp_path,
        file_stem="fake-chart",
    )
    prs = ctx.new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.add_textbox(0, 0, 1000000, 500000).text_frame.text = "chart_region"
    manifest = ctx.save_artifacts(prs, evidence_refs_used=["chart_region"])

    report = validate_artifacts(manifest)

    assert report.status == "failed_validation"
    assert any(f.code == "CHART_REF_WITHOUT_NATIVE_CHART" for f in report.findings)


def test_renderer_guard_allows_visible_numbers_for_artifact_warning() -> None:
    source = '''
from lobster_runtime.smart_report_pptx import Deck
def render(ctx):
    deck = Deck.from_context(ctx)
    slide = deck.add_slide(role="content")
    slide.add_free_text("營收成長 12.5%", box=(0.5, 0.5, 3.0, 1.0))
    return deck.save()
'''.strip()
    validate_renderer_source(source)


def test_renderer_guard_allows_visible_numbers_in_text_aliases() -> None:
    source = '''
from lobster_runtime.smart_report_pptx import Deck
def render(ctx):
    deck = Deck.from_context(ctx)
    slide = deck.add_slide(role="content")
    slide.add_text("營收成長 12.5%", box=(0.5, 0.5, 3.0, 1.0))
    slide.add_bullet_list(["市占率 30%"], box=(0.5, 1.6, 3.0, 1.0))
    return deck.save()
'''.strip()
    validate_renderer_source(source)


def test_renderer_guard_allows_visible_numbers_in_native_text_assignment() -> None:
    source = '''
from pptx.util import Inches
def render(ctx):
    prs = ctx.new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3.0), Inches(1.0))
    box.text_frame.text = "營收成長 12.5%"
    return ctx.save_artifacts(prs)
'''.strip()
    validate_renderer_source(source)


def test_renderer_guard_rejects_external_workbook_reads() -> None:
    source = '''
from openpyxl import load_workbook
def render(ctx):
    load_workbook("anything.xlsx")
'''.strip()
    with pytest.raises(ValueError, match="import is not allowed"):
        validate_renderer_source(source)


def test_renderer_guard_allows_common_python_pptx_font_helpers() -> None:
    source = '''
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
def render(ctx):
    prs = ctx.new_presentation()
    return ctx.save_artifacts(prs)
'''.strip()
    validate_renderer_source(source)


def test_full_universal_pipeline_reaches_artifacts_with_agent_stages(tmp_path) -> None:
    workbook_path = tmp_path / "demo.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Region"
    sheet.append(["區域", "營收"])
    sheet.append(["北區", 600000])
    sheet.append(["中區", 350000])
    sheet.append(["南區", 300000])
    workbook.save(workbook_path)

    class ScriptedUniversalPipeline(UniversalPresentationPipeline):
        def _run_stage(self, stage_name, output_model, system_prompt, context):  # type: ignore[no-untyped-def]
            if output_model is DataIntelligenceReport:
                return DataIntelligenceReport(
                    status="passed",
                    data_structure_summary="Workbook contains a region table with numeric revenue.",
                    semantic_notes=["Region is categorical and revenue is numeric."],
                    relationship_hypotheses=[],
                    quality_findings=[],
                    usable_data_notes=["The table can support a summary presentation."],
                )
            if output_model is AnalysisFeasibilityPlan:
                return AnalysisFeasibilityPlan(
                    status="passed",
                    accepted_analyses=[{
                        "analysis_id": "analysis_region_revenue",
                        "question": "How does revenue compare by region?",
                        "method": "Sum numeric revenue columns and compare categories.",
                        "required_columns": ["區域", "營收"],
                        "feasible": True,
                        "rationale": "The workbook has visible category and numeric values.",
                        "recommended_visual": "column chart",
                    }],
                    rejected_analyses=[],
                    known_limits=[],
                    questions_for_user=[],
                )
            if output_model is EvidencePacketV2:
                return EvidencePacketV2.model_validate({
                    "packet_id": "agent-evidence-test",
                    "metrics": [{
                        "metric_id": "metric_1_營收",
                        "label": "營收",
                        "value": 1250000,
                        "unit": "元",
                        "display_format": "currency",
                        "source_refs": ["demo.xlsx#Region!營收"],
                        "calculation": "agent selected revenue summary from workbook profile",
                    }],
                    "charts": [{
                        "chart_id": "chart_1_營收",
                        "title": "區域營收",
                        "chart_type": "column",
                        "categories": ["北區", "中區", "南區"],
                        "series": [{"name": "營收", "values": [600000, 350000, 300000]}],
                        "metric_refs": ["metric_1_營收"],
                    }],
                    "claims": [{
                        "claim_id": "claim_1_summary",
                        "text": "本期營收為 {{metric_1_營收}}，後續應聚焦主要區域。",
                        "metric_refs": ["metric_1_營收"],
                        "chart_refs": ["chart_1_營收"],
                        "source_refs": ["demo.xlsx#Region"],
                    }],
                    "tables": [{
                        "table_id": "table_1_營收",
                        "title": "區域營收表",
                        "headers": ["區域", "營收"],
                        "rows": [["北區", 600000], ["中區", 350000], ["南區", 300000]],
                        "source_refs": ["demo.xlsx#Region"],
                    }],
                })
            if output_model is VerifiedAnalysisNarrative:
                return VerifiedAnalysisNarrative(
                    status="passed",
                    insight_summaries=["Use the generated evidence packet for all numeric content."],
                    caveats=[],
                    evidence_usage_notes=["Claims and charts must reference evidence IDs."],
                )
            if output_model is BlueprintStageOutput:
                return BlueprintStageOutput(
                    status="passed",
                    blueprint={
                        "blueprint_version": "presentation-blueprint-v1",
                        "title": "營運表現摘要",
                        "slides": [{
                            "slide_id": "s1",
                            "role": "cover",
                            "intent": "Summarize evidence-backed result.",
                            "layout_strategy": "claim led title slide",
                            "elements": [
                                {"element_id": "e1", "type": "title", "box": {"x": 0.7, "y": 0.7, "w": 10.5, "h": 0.8}, "text": "營運表現摘要"},
                                {"element_id": "e2", "type": "claim", "box": {"x": 0.9, "y": 1.8, "w": 7.5, "h": 1.2}, "claim_ref": "claim_1_summary"},
                                {"element_id": "e3", "type": "chart", "box": {"x": 0.9, "y": 3.0, "w": 6.0, "h": 3.0}, "chart_ref": "chart_1_營收"},
                            ],
                        }],
                    },
                    design_notes=["Use template if present."],
                )
            raise AssertionError(f"unexpected stage model: {output_model}")

    manifest = ScriptedUniversalPipeline(ScriptedRendererModel(source=PIPELINE_RENDERER_SOURCE)).run(
        prompt="產生一份區域營收簡報",
        data_paths=[workbook_path],
        output_dir=tmp_path,
    )

    assert manifest.status == "final"
    assert manifest.data_report.status == "passed"
    assert manifest.feasibility_plan.accepted_analyses
    assert manifest.pptx_path.endswith(".pptx")

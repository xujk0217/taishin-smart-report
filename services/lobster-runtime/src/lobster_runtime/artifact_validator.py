"""Validators that inspect the actual generated PPTX and XLSX files."""

from __future__ import annotations

from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation

from .presentation_contracts import ArtifactFinding, ArtifactValidationReport, RenderArtifactManifest
from .smart_report_pptx import CONTENT_NUMBER_RE


def validate_artifacts(manifest: RenderArtifactManifest) -> ArtifactValidationReport:
    findings: list[ArtifactFinding] = []
    slide_count = 0
    chart_count = manifest.chart_count
    table_count = manifest.table_count

    try:
        presentation = Presentation(manifest.pptx_path)
        slide_count = len(presentation.slides)
        actual_chart_count = 0
        if slide_count != manifest.slide_count:
            findings.append(_finding("blocking", "SLIDE_COUNT_MISMATCH", "PPTX slide count does not match render manifest"))
        for index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if getattr(shape, "has_chart", False):
                    actual_chart_count += 1
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text_frame.text
                if "placeholder" in text.lower() or "test data" in text.lower():
                    findings.append(_finding("blocking", "PLACEHOLDER_TEXT", f"Slide {index} contains placeholder/test text"))
                for match in CONTENT_NUMBER_RE.finditer(text):
                    value = match.group(0).strip()
                    if value and value not in _allowed_rendered_numbers(manifest):
                        findings.append(_finding(
                            "warning",
                            "POSSIBLE_UNTRACED_NUMBER",
                            f"Slide {index} contains a number-like text fragment that must be checked: {value}",
                        ))
        chart_count = actual_chart_count
        if manifest.chart_count > 0 and actual_chart_count < manifest.chart_count:
            findings.append(_finding(
                "blocking",
                "NATIVE_CHART_COUNT_MISMATCH",
                "Render manifest reports charts, but the PPTX contains fewer native editable chart objects",
            ))
        if actual_chart_count == 0 and any("chart" in ref.lower() for ref in manifest.evidence_refs_used):
            findings.append(_finding(
                "blocking",
                "CHART_REF_WITHOUT_NATIVE_CHART",
                "Chart evidence was used, but the PPTX contains no native editable chart object",
            ))
    except Exception as error:
        findings.append(_finding("blocking", "PPTX_OPEN_FAILED", f"Generated PPTX could not be opened: {type(error).__name__}"))

    try:
        workbook = load_workbook(manifest.xlsx_path, read_only=True, data_only=True)
        workbook.close()
    except Exception as error:
        findings.append(_finding("blocking", "XLSX_OPEN_FAILED", f"Generated XLSX could not be opened: {type(error).__name__}"))

    if not Path(manifest.pptx_path).exists() or Path(manifest.pptx_path).stat().st_size == 0:
        findings.append(_finding("blocking", "PPTX_EMPTY", "Generated PPTX is missing or empty"))
    if not Path(manifest.xlsx_path).exists() or Path(manifest.xlsx_path).stat().st_size == 0:
        findings.append(_finding("blocking", "XLSX_EMPTY", "Generated XLSX is missing or empty"))

    status = "failed_validation" if any(f.severity == "blocking" for f in findings) else "final"
    return ArtifactValidationReport(
        status=status,
        findings=findings,
        checked_slide_count=slide_count,
        checked_chart_count=chart_count,
        checked_table_count=table_count,
    )


def _allowed_rendered_numbers(manifest: RenderArtifactManifest) -> set[str]:
    # This first pass is intentionally conservative. Numbers from metric cards
    # and claims are allowed by SDK construction, so artifact extraction emits
    # warnings rather than blocking when it sees number-like text.
    return {str(manifest.slide_count), str(manifest.chart_count), str(manifest.table_count)}


def _finding(severity: str, code: str, message: str) -> ArtifactFinding:
    return ArtifactFinding(
        severity=severity,  # type: ignore[arg-type]
        code=code,
        message=message,
        origin_stage="artifact-validation",
    )

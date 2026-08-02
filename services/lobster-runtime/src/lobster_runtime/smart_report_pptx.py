"""Controlled SDK exposed to generated python-pptx renderer programs."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from openpyxl import Workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .presentation_contracts import (
    ChartData,
    ClaimRecord,
    EvidenceMetric,
    EvidencePacketV2,
    RenderArtifactManifest,
    TableData,
)


CONTENT_NUMBER_RE = re.compile(r"(?<![#A-Za-z0-9])(?:[$NTD€¥]\s*)?\d+(?:,\d{3})*(?:\.\d+)?\s*(?:%|％|元|萬|億|千|人|件|名|次|kg|km|分|秒)?")


class EvidenceIndex(dict[str, Any]):
    """Dict-like lookup that also behaves like a list of evidence records."""

    def __getitem__(self, key: str | int | slice) -> Any:
        if isinstance(key, int):
            return list(self.values())[key]
        if isinstance(key, slice):
            return list(self.values())[key]
        return super().__getitem__(key)

    def __iter__(self):
        return iter(self.values())


class RenderingContext:
    """Immutable render input available to generated code as ctx."""

    def __init__(
        self,
        *,
        evidence: EvidencePacketV2,
        output_dir: str | Path,
        template_path: str | Path | None = None,
        file_stem: str = "agent-generated-presentation",
    ) -> None:
        self.evidence = evidence
        self.output_dir = Path(output_dir)
        self.template_path = Path(template_path) if template_path else None
        self.file_stem = file_stem
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.last_deck: Deck | None = None
        self.last_presentation: Presentation | None = None
        self.last_workbook: Workbook | None = None
        self.metrics = EvidenceIndex({metric.metric_id: metric for metric in evidence.metrics})
        self.charts = EvidenceIndex({chart.chart_id: chart for chart in evidence.charts})
        self.claims = EvidenceIndex({claim.claim_id: claim for claim in evidence.claims})
        self.tables = EvidenceIndex({table.table_id: table for table in evidence.tables})

    @property
    def output_pptx_path(self) -> str:
        return str(self.output_dir / f"{self.file_stem}.pptx")

    @property
    def output_xlsx_path(self) -> str:
        return str(self.output_dir / f"{self.file_stem}.xlsx")

    @property
    def pptx_path(self) -> str:
        return self.output_pptx_path

    @property
    def xlsx_path(self) -> str:
        return self.output_xlsx_path

    @property
    def template_file(self) -> str | None:
        return str(self.template_path) if self.template_path else None

    def metric(self, metric_id: str) -> EvidenceMetric:
        return self.metrics[metric_id]

    def chart(self, chart_id: str) -> ChartData:
        return self.charts[chart_id]

    def claim(self, claim_id: str) -> ClaimRecord:
        return self.claims[claim_id]

    def table(self, table_id: str) -> TableData:
        return self.tables[table_id]

    def new_presentation(
        self,
        preserve_template_slides: bool = False,
        *,
        keep_template_slides: bool | None = None,
        use_template_components: bool | None = None,
    ) -> Presentation:
        presentation = Presentation(str(self.template_path)) if self.template_path else Presentation()
        preserve = preserve_template_slides or bool(keep_template_slides) or bool(use_template_components)
        if self.template_path and not preserve:
            _remove_existing_slides(presentation)
        self.last_presentation = presentation
        return presentation

    def remove_slide(self, presentation: Presentation, index: int) -> None:
        slide_id_list = presentation.slides._sldIdLst
        slide_id = list(slide_id_list)[index]
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)

    def clear_slide_text(self, slide: Any) -> None:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.clear()

    def new_workbook(self) -> Workbook:
        workbook = Workbook()
        workbook.remove(workbook.active)
        self.last_workbook = workbook
        return workbook

    def save_artifacts(
        self,
        presentation: Presentation | None = None,
        workbook: Workbook | None = None,
        *,
        chart_count: int = 0,
        table_count: int = 0,
        evidence_refs_used: list[str] | set[str] | tuple[str, ...] | None = None,
        **kwargs: Any,
    ) -> RenderArtifactManifest:
        chart_count = int(kwargs.get("charts_count", chart_count) or 0)
        table_count = int(kwargs.get("tables_count", table_count) or 0)
        evidence_refs_used = evidence_refs_used or kwargs.get("evidence_refs") or kwargs.get("refs_used")
        pptx = presentation or self.last_presentation
        if pptx is None:
            raise ValueError("presentation is required")
        xlsx = workbook or self.last_workbook or self.new_workbook()
        if not xlsx.sheetnames:
            xlsx.create_sheet("Manifest")
            xlsx["Manifest"].append(["message"])
            xlsx["Manifest"].append(["No chart or table data was rendered."])
        pptx.save(self.output_pptx_path)
        xlsx.save(self.output_xlsx_path)
        return RenderArtifactManifest(
            pptx_path=self.output_pptx_path,
            xlsx_path=self.output_xlsx_path,
            slide_count=len(pptx.slides),
            chart_count=chart_count,
            table_count=table_count,
            evidence_refs_used=sorted(evidence_refs_used or []),
        )


class Deck:
    """Small safe wrapper around python-pptx plus a synchronized XLSX workbook."""

    def __init__(self, ctx: RenderingContext) -> None:
        self.ctx = ctx
        ctx.last_deck = self
        self.presentation = ctx.new_presentation()
        self.workbook = ctx.new_workbook()
        self._chart_count = 0
        self._table_count = 0
        self._evidence_refs_used: set[str] = set()

    @classmethod
    def from_context(cls, ctx: RenderingContext) -> "Deck":
        return cls(ctx)

    def use_template(self, *_args: Any, **_kwargs: Any) -> "Deck":
        """Compatibility hook: templates are applied when Deck is constructed."""
        return self

    def add_slide(
        self,
        role: str = "content",
        layout_index: int | None = None,
        *,
        layout_name: str | None = None,
        **_: Any,
    ) -> "Slide":
        layouts = self.presentation.slide_layouts
        index = (
            layout_index
            if layout_index is not None
            else self._layout_for_role(layouts, role, layout_name=layout_name)
        )
        index = max(0, min(index, len(layouts) - 1))
        return Slide(self, self.presentation.slides.add_slide(layouts[index]), role=role)

    def save(self) -> RenderArtifactManifest:
        if not self.workbook.sheetnames:
            self.workbook.create_sheet("Manifest")
            self.workbook["Manifest"].append(["message"])
            self.workbook["Manifest"].append(["No chart or table data was rendered."])
        return self.ctx.save_artifacts(
            self.presentation,
            self.workbook,
            chart_count=self._chart_count,
            table_count=self._table_count,
            evidence_refs_used=self._evidence_refs_used,
        )

    @classmethod
    def _layout_for_role(cls, layouts: Any, role: str, *, layout_name: str | None = None) -> int:
        count = len(layouts)
        if count <= 1:
            return 0
        if layout_name:
            exact = cls._find_layout(layouts, equals=[layout_name])
            if exact is not None:
                return exact
        if role in {"cover", "section", "back-cover"}:
            preferred = {
                "cover": [["2_", "標題投影片"], ["cover"], ["title"]],
                "section": [["章節"], ["section"]],
                "back-cover": [["3_", "標題投影片"], ["back"], ["closing"]],
            }[role]
            found = cls._find_layout(layouts, contains_all=preferred)
            if found is not None:
                return found
            return 0
        found = cls._find_layout(layouts, contains_all=[["標題", "內容"], ["title", "content"]])
        if found is not None:
            return found
        return min(1, count - 1)

    @staticmethod
    def _find_layout(
        layouts: Any,
        *,
        equals: list[str] | None = None,
        contains_all: list[list[str]] | None = None,
    ) -> int | None:
        equals_normalized = {_normalize_layout_name(value) for value in equals or []}
        for index, layout in enumerate(layouts):
            name = _normalize_layout_name(getattr(layout, "name", ""))
            if name in equals_normalized:
                return index
        for keywords in contains_all or []:
            normalized_keywords = [_normalize_layout_name(value) for value in keywords]
            for index, layout in enumerate(layouts):
                name = _normalize_layout_name(getattr(layout, "name", ""))
                if all(keyword in name for keyword in normalized_keywords):
                    return index
        return None


class Slide:
    def __init__(self, deck: Deck, slide: Any, *, role: str = "content") -> None:
        self.deck = deck
        self.slide = slide
        self.role = role

    def add_title(self, content: str | None = None, *, text: str | None = None, box: Any = None, **_: Any) -> None:
        rendered_text = _coalesce_text(content, text)
        _reject_untraced_content_numbers(rendered_text)
        shape = self._template_text_shape("title")
        if shape is None:
            shape = self.slide.shapes.add_textbox(*_box(box))
        _write_text(shape.text_frame, rendered_text, size=24, bold=True, alignment=PP_ALIGN.CENTER if self.role in {"cover", "section", "back-cover"} else PP_ALIGN.LEFT)

    def add_subtitle(self, content: str | None = None, *, text: str | None = None, box: Any = None, **kwargs: Any) -> None:
        self.add_free_text(content, text=text, box=box, **kwargs)

    def add_free_text(self, content: str | None = None, *, text: str | None = None, box: Any = None, **_: Any) -> None:
        rendered_text = _coalesce_text(content, text)
        _reject_untraced_content_numbers(rendered_text)
        shape = self.slide.shapes.add_textbox(*_box(box))
        _write_text(shape.text_frame, rendered_text, size=13)

    def add_text(self, content: str | None = None, *, text: str | None = None, box: Any = None, **kwargs: Any) -> None:
        self.add_free_text(content, text=text, box=box, **kwargs)

    def add_summary(self, content: str | None = None, *, text: str | None = None, box: Any = None, **kwargs: Any) -> None:
        self.add_free_text(content, text=text, box=box)

    def add_bullet_list(self, items: list[str] | tuple[str, ...] | str | None = None, *, box: Any = None, **_: Any) -> None:
        if isinstance(items, str):
            bullet_items = [items]
        else:
            bullet_items = list(items or [])
        for item in bullet_items:
            _reject_untraced_content_numbers(item)
        shape = self._template_text_shape("body") if box is None else None
        if shape is None:
            shape = self.slide.shapes.add_textbox(*_box(box))
        frame = shape.text_frame
        frame.clear()
        for index, item in enumerate(bullet_items or [""]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            paragraph.font.size = Pt(13)

    def add_metric_card(self, metric_ref: str | None = None, *, id: str | None = None, ref: str | None = None, label: str | None = None, box: Any = None, **_: Any) -> None:
        metric = self._metric(_coalesce_ref(metric_ref, id, ref, label="metric_ref"))
        self.deck._evidence_refs_used.add(metric.metric_id)
        shape = self.slide.shapes.add_textbox(*_box(box))
        frame = shape.text_frame
        frame.clear()
        frame.paragraphs[0].text = label or metric.label
        frame.paragraphs[0].font.size = Pt(11)
        paragraph = frame.add_paragraph()
        paragraph.text = _format_metric(metric)
        paragraph.font.size = Pt(22)
        paragraph.font.bold = True

    def add_claim(self, claim_ref: str | None = None, *, id: str | None = None, ref: str | None = None, box: Any = None, **_: Any) -> None:
        claim = self._claim(_coalesce_ref(claim_ref, id, ref, label="claim_ref"))
        self.deck._evidence_refs_used.add(claim.claim_id)
        self.deck._evidence_refs_used.update(claim.metric_refs)
        shape = self.slide.shapes.add_textbox(*_box(box))
        _write_text(shape.text_frame, _render_claim_text(claim, self.deck.ctx.metrics), size=14)

    def add_table(self, table_ref: str | None = None, *, id: str | None = None, ref: str | None = None, box: Any = None, **_: Any) -> None:
        table = self._table(_coalesce_ref(table_ref, id, ref, label="table_ref"))
        self.deck._evidence_refs_used.add(table.table_id)
        rows = max(1, len(table.rows) + 1)
        cols = len(table.headers)
        ppt_table = self.slide.shapes.add_table(rows, cols, *_box(box)).table
        for col_index, header in enumerate(table.headers):
            ppt_table.cell(0, col_index).text = header
        for row_index, row in enumerate(table.rows, start=1):
            for col_index, value in enumerate(row[:cols]):
                ppt_table.cell(row_index, col_index).text = str(value)
        sheet = self.deck.workbook.create_sheet(_safe_sheet_name(table.table_id))
        sheet.append(table.headers)
        for row in table.rows:
            sheet.append(row)
        self.deck._table_count += 1

    def add_chart(self, chart_ref: str | None = None, *, id: str | None = None, ref: str | None = None, box: Any = None, **_: Any) -> None:
        chart = self._chart(_coalesce_ref(chart_ref, id, ref, label="chart_ref"))
        self.deck._evidence_refs_used.add(chart.chart_id)
        self.deck._evidence_refs_used.update(chart.metric_refs)
        chart_data = CategoryChartData()
        chart_data.categories = chart.categories
        for series in chart.series:
            chart_data.add_series(series.name, series.values)
        self.slide.shapes.add_chart(_chart_type(chart), *_box(box), chart_data)
        sheet = self.deck.workbook.create_sheet(_safe_sheet_name(chart.chart_id))
        sheet.append(["category", *[series.name for series in chart.series]])
        for index, category in enumerate(chart.categories):
            sheet.append([category, *[series.values[index] for series in chart.series]])
        self.deck._chart_count += 1

    def _metric(self, metric_ref: str) -> EvidenceMetric:
        try:
            return self.deck.ctx.metrics[metric_ref]
        except KeyError as error:
            raise ValueError(f"unknown metric_ref: {metric_ref}") from error

    def _chart(self, chart_ref: str) -> ChartData:
        try:
            return self.deck.ctx.charts[chart_ref]
        except KeyError as error:
            raise ValueError(f"unknown chart_ref: {chart_ref}") from error

    def _claim(self, claim_ref: str) -> ClaimRecord:
        try:
            return self.deck.ctx.claims[claim_ref]
        except KeyError as error:
            raise ValueError(f"unknown claim_ref: {claim_ref}") from error

    def _table(self, table_ref: str) -> TableData:
        try:
            return self.deck.ctx.tables[table_ref]
        except KeyError as error:
            raise ValueError(f"unknown table_ref: {table_ref}") from error

    def _template_text_shape(self, kind: str) -> Any | None:
        candidates: list[tuple[int, float, Any]] = []
        for shape in self.slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            name = _normalize_layout_name(getattr(shape, "name", ""))
            if any(token in name for token in ("編號", "頁碼", "footer", "date")):
                continue
            top = float(shape.top) / 914400
            left = float(shape.left) / 914400
            area = float(shape.width * shape.height)
            if kind == "title" and ("標題" in name or "title" in name):
                candidates.append((0, top + left / 100, shape))
            elif kind == "body" and any(token in name for token in ("內容", "文字", "content", "body")):
                candidates.append((0, -area, shape))
        if candidates:
            candidates.sort(key=lambda item: (item[0], item[1]))
            return candidates[0][2]
        return None


def _box(box: Any) -> tuple[Any, Any, Any, Any]:
    x, y, w, h = _coerce_box_values(box)
    return Inches(x), Inches(y), Inches(w), Inches(h)


def _normalize_layout_name(value: str) -> str:
    return str(value).strip().lower().replace(" ", "")


def _write_text(frame: Any, text: str, *, size: int, bold: bool = False, alignment: Any | None = None) -> None:
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    if alignment is not None:
        paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold


def _coerce_box_values(box: Any) -> tuple[float, float, float, float]:
    default = (0.7, 0.7, 5.5, 1.0)
    if isinstance(box, dict):
        raw_values = [box.get("x"), box.get("y"), box.get("w"), box.get("h")]
    elif hasattr(box, "x") and hasattr(box, "y") and hasattr(box, "w") and hasattr(box, "h"):
        raw_values = [box.x, box.y, box.w, box.h]
    elif isinstance(box, list | tuple) and len(box) >= 4:
        raw_values = list(box[:4])
    else:
        return default
    try:
        x, y, w, h = (float(value) for value in raw_values)
    except (TypeError, ValueError):
        return default
    if w <= 0 or h <= 0:
        return default
    return (
        max(0.0, min(x, 12.8)),
        max(0.0, min(y, 7.1)),
        max(0.1, min(w, 13.3)),
        max(0.1, min(h, 7.5)),
    )


def _coalesce_text(content: str | None, text: str | None) -> str:
    if text is not None:
        return text
    if content is not None:
        return content
    raise ValueError("text content is required")


def _coalesce_ref(primary: str | None, id_value: str | None, ref: str | None, *, label: str) -> str:
    value = primary or id_value or ref
    if value:
        return value
    raise ValueError(f"{label} is required")


def _remove_existing_slides(presentation: Presentation) -> None:
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _format_metric(metric: EvidenceMetric) -> str:
    if metric.display_format == "percent":
        return f"{metric.value:.1%}" if abs(metric.value) <= 1 else f"{metric.value:.1f}%"
    if metric.display_format == "currency":
        return f"{metric.value:,.0f}{metric.unit}"
    if metric.display_format == "integer":
        return f"{metric.value:,.0f}{metric.unit}"
    return f"{metric.value:,.2f}{metric.unit}".rstrip("0").rstrip(".")


def _render_claim_text(claim: ClaimRecord, metrics: dict[str, EvidenceMetric]) -> str:
    text = claim.text
    for metric_ref in claim.metric_refs:
        if metric_ref in metrics:
            text = text.replace(f"{{{{{metric_ref}}}}}", _format_metric(metrics[metric_ref]))
    unresolved_tokens = re.findall(r"\{\{[^}]+\}\}", text)
    if unresolved_tokens:
        raise ValueError(f"claim contains unresolved metric placeholders: {', '.join(unresolved_tokens)}")
    _reject_untraced_content_numbers(text, allow=True)
    return text


def _reject_untraced_content_numbers(text: str, *, allow: bool = False) -> None:
    if allow:
        return
    matches = [match.group(0).strip() for match in CONTENT_NUMBER_RE.finditer(text)]
    material = [value for value in matches if _is_explicit_material_number(value)]
    if material:
        raise ValueError(f"content text contains untraced numbers: {', '.join(material)}")


def _is_explicit_material_number(value: str) -> bool:
    return bool(re.search(r"[$NTD€¥%％元萬億千人件名次kgkm分秒]", value))


def _chart_type(chart: ChartData) -> Any:
    return {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }[chart.chart_type]


def _safe_sheet_name(value: str) -> str:
    cleaned = re.sub(r"[\[\]:*?/\\]", "_", value)
    return cleaned[:31] or "Sheet"

"""PowerPoint template profiling for agent presentation generation."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from .presentation_contracts import Box, TemplateProfile


def analyze_template(template_path: str | Path | None) -> TemplateProfile:
    """Read a user-provided PPTX template, or return a default profile."""
    if template_path is None:
        return TemplateProfile(
            source="default",
            layouts=["default-title", "default-content"],
            theme_fonts=["Aptos", "Arial"],
            theme_colors=["#A50034", "#FFFFFF", "#404040"],
            fixed_regions=[Box(x=0, y=7.0, w=13.333, h=0.5)],
        )
    path = Path(template_path)
    presentation = Presentation(str(path))
    layouts = [layout.name or f"layout-{index}" for index, layout in enumerate(presentation.slide_layouts)]
    width = presentation.slide_width / 914400
    height = presentation.slide_height / 914400
    return TemplateProfile(
        template_path=str(path),
        source="uploaded",
        slide_width=width,
        slide_height=height,
        layouts=layouts,
        theme_fonts=[],
        theme_colors=[],
        fixed_regions=[Box(x=0, y=max(0, height - 0.5), w=width, h=0.5)],
        sample_slides=_sample_slides(presentation),
    )


def _sample_slides(presentation: Presentation) -> list[dict[str, object]]:
    slides: list[dict[str, object]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        shapes: list[dict[str, object]] = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = ""
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text[:120]
            shapes.append({
                "index": shape_index,
                "name": getattr(shape, "name", ""),
                "shape_type": str(getattr(shape, "shape_type", "")),
                "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
                "text": text,
                "x": round(shape.left / 914400, 3),
                "y": round(shape.top / 914400, 3),
                "w": round(shape.width / 914400, 3),
                "h": round(shape.height / 914400, 3),
            })
        slides.append({
            "slide_index": slide_index,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "",
            "shapes": shapes,
        })
    return slides

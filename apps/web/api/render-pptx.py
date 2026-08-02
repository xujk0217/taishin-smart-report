"""
Vercel Serverless Function (Python) — renders a PPTX from a user-uploaded
template + slide spec JSON.

POST /api/render-pptx
  Content-Type: multipart/form-data
  Fields:
    template: .pptx file (the user's brand template)
    spec: JSON string (the SlideSpec array)
    data: JSON string (computed metrics/charts for populating content)

Returns: application/vnd.openxmlformats-officedocument.presentationml.presentation
"""
from http.server import BaseHTTPRequestHandler
import json
import io
import base64
import tempfile
import os

# python-pptx must be available; Vercel installs from requirements.txt
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError as e:
    Presentation = None
    IMPORT_ERROR = str(e)


class handler(BaseHTTPRequestHandler):
    def do_OPTIONS(self):
        self.send_response(200)
        self._cors()
        self.end_headers()

    def do_POST(self):
        if Presentation is None:
            self._error(500, f'python-pptx not available: {IMPORT_ERROR}')
            return

        content_type = self.headers.get('Content-Type', '')

        try:
            if 'multipart/form-data' in content_type:
                result = self._handle_multipart()
            elif 'application/json' in content_type:
                result = self._handle_json()
            else:
                self._error(400, 'Expected multipart/form-data or application/json')
                return

            if result is None:
                return  # error already sent

            pptx_bytes = result
            self.send_response(200)
            self._cors()
            self.send_header('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            self.send_header('Content-Disposition', 'attachment; filename="report.pptx"')
            self.send_header('Content-Length', str(len(pptx_bytes)))
            self.end_headers()
            self.wfile.write(pptx_bytes)

        except Exception as e:
            self._error(500, f'Render failed: {str(e)[:200]}')

    def _handle_json(self):
        """For JSON-only requests (template sent as base64 or omitted)."""
        length = int(self.headers.get('Content-Length', 0))
        body = json.loads(self.rfile.read(length))

        template_b64 = body.get('template')  # base64-encoded .pptx
        spec = body.get('spec', [])
        data = body.get('data', {})

        if template_b64:
            template_bytes = base64.b64decode(template_b64)
            prs = Presentation(io.BytesIO(template_bytes))
        else:
            # No template: create a blank 16:9 presentation
            prs = Presentation()
            prs.slide_width = Emu(12192000)  # 16:9
            prs.slide_height = Emu(6858000)

        return self._render(prs, spec, data)

    def _handle_multipart(self):
        """For multipart form data (template as file upload)."""
        import cgi
        form = cgi.FieldStorage(
            fp=self.rfile,
            headers=self.headers,
            environ={'REQUEST_METHOD': 'POST', 'CONTENT_TYPE': self.headers['Content-Type']}
        )

        template_field = form['template'] if 'template' in form else None
        spec_field = form.getvalue('spec', '[]')
        data_field = form.getvalue('data', '{}')

        spec = json.loads(spec_field) if isinstance(spec_field, str) else spec_field
        data = json.loads(data_field) if isinstance(data_field, str) else data_field

        if template_field and template_field.file:
            prs = Presentation(template_field.file)
        else:
            prs = Presentation()
            prs.slide_width = Emu(12192000)
            prs.slide_height = Emu(6858000)

        return self._render(prs, spec, data)

    def _render(self, prs, spec, data):
        """Render exact approved pages from the template's reusable layouts.

        Template sample slides are deliberately not overwritten: they can contain
        example copy or fake values. New slides inherit the selected template
        layout, master, theme, and placeholders instead.
        """
        layouts = list(prs.slide_layouts)
        if not layouts:
            raise ValueError('Template contains no reusable slide layouts')

        # A template is often an example deck.  Its existing slides can contain
        # demo copy or invented values, so keep only the reusable layouts/master
        # and replace the deck with exactly the approved plan pages.
        self._remove_existing_slides(prs)
        for slide_spec in spec:
            chosen_layout = self._select_layout(layouts, slide_spec.get('layout', 'content'))
            slide = prs.slides.add_slide(chosen_layout)
            self._render_elements(slide, slide_spec.get('elements', []), slide_spec, data, prs)

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    @staticmethod
    def _remove_existing_slides(prs):
        """Remove slides while retaining the template's layouts and theme.

        python-pptx has no public delete-slide API.  Removing the slide id and
        its relationship is the supported low-level pattern for this operation;
        it leaves slide masters and layouts intact for subsequent add_slide().
        """
        while prs.slides:
            slide_id = prs.slides._sldIdLst[0]
            prs.part.drop_rel(slide_id.rId)
            prs.slides._sldIdLst.remove(slide_id)

    def _render_elements(self, slide, elements, slide_spec, data, prs):
        """Add elements to a slide."""
        sw = prs.slide_width
        sh = prs.slide_height
        margin = Inches(0.6)
        content_w = sw - margin * 2
        y_cursor = Inches(0.8)
        layout_type = slide_spec.get('layout', 'content')
        is_centered = layout_type in ('cover', 'section_title', 'backcover')
        chart_element = next((el for el in elements if el.get('type') == 'chart'), None)
        chart_box = None
        if chart_element:
            # A chart must occupy the template's intended content frame.  Do
            # not layer it over an inherited body placeholder or title region.
            body_placeholder = self._find_placeholder(slide, (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT))
            if body_placeholder is not None:
                chart_box = (body_placeholder.left, body_placeholder.top, body_placeholder.width, body_placeholder.height)
                self._remove_shape(body_placeholder)

        # Prefer inherited title/body placeholders. This preserves font, size,
        # alignment and other template-defined typography. Only add a textbox
        # when the selected layout genuinely has no matching placeholder.
        title_element = next((el for el in elements if el.get('type') == 'title' and el.get('content')), None)
        if title_element and self._write_placeholder(slide, title_element['content'], 'title'):
            elements = [el for el in elements if el is not title_element]
            inherited_title = slide.shapes.title
            if inherited_title is not None:
                y_cursor = max(y_cursor, inherited_title.top + inherited_title.height + Inches(0.05))
        body_element = None if chart_element else next((el for el in elements if el.get('type') in ('subtitle', 'text_block', 'bullet_list') and (el.get('content') or el.get('items'))), None)
        if body_element:
            body_text = body_element.get('content') or '\n'.join(f'• {item}' for item in body_element.get('items', []))
            if self._write_placeholder(slide, body_text, 'body'):
                elements = [el for el in elements if el is not body_element]

        for el in elements:
            el_type = el.get('type', '')
            content = el.get('content', '')
            size = el.get('size', 'medium')

            # The inherited body frame is dedicated to the chart.  A chart
            # slide gets its narrative claim from title/subtitle; placing the
            # full bullet list into the same frame causes overlap and makes
            # the chart hard to read.
            if chart_element and el_type in ('heading', 'text_block', 'bullet_list', 'kpi_block', 'comparison', 'table', 'insight'):
                continue

            if el_type == 'title':
                h = Inches(1.2) if is_centered else Inches(0.7)
                y = Inches(2.5) if is_centered else y_cursor
                tf = slide.shapes.add_textbox(margin, y, content_w, h).text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(36 if layout_type == 'cover' else 28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
                if is_centered:
                    p.alignment = PP_ALIGN.CENTER
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                if not is_centered:
                    y_cursor += h + Inches(0.1)

            elif el_type == 'subtitle':
                h = Inches(0.6)
                y = Inches(3.8) if is_centered else y_cursor
                tf = slide.shapes.add_textbox(margin, y, content_w, h).text_frame
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
                if is_centered:
                    p.alignment = PP_ALIGN.CENTER
                if not is_centered:
                    y_cursor += h + Inches(0.1)

            elif el_type == 'heading':
                h = Inches(0.5)
                tf = slide.shapes.add_textbox(margin, y_cursor, content_w, h).text_frame
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
                y_cursor += h + Inches(0.15)

            elif el_type == 'text_block':
                h = Inches(1.2) if size in ('large', 'full') else Inches(0.8)
                tf = slide.shapes.add_textbox(margin, y_cursor, content_w, h).text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
                y_cursor += h + Inches(0.1)

            elif el_type == 'bullet_list':
                items = el.get('items', [])
                h = Inches(0.3 * len(items) + 0.2)
                tf = slide.shapes.add_textbox(margin, y_cursor, content_w, h).text_frame
                tf.word_wrap = True
                for j, item in enumerate(items):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = f'• {item}'
                    p.font.size = Pt(11)
                    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
                    p.space_before = Pt(4)
                y_cursor += h + Inches(0.1)

            elif el_type == 'kpi_block':
                metrics = el.get('metrics', [])
                for m in metrics:
                    h = Inches(0.5)
                    tf = slide.shapes.add_textbox(margin, y_cursor, Inches(3), h).text_frame
                    p = tf.paragraphs[0]
                    p.text = m.get('value', '')
                    p.font.size = Pt(22)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)

                    label_tf = slide.shapes.add_textbox(margin + Inches(3.2), y_cursor, Inches(4), h).text_frame
                    lp = label_tf.paragraphs[0]
                    label_parts = [m.get('label', '')]
                    if m.get('rank'):
                        label_parts.append(f'#{m["rank"]}')
                    if m.get('trend'):
                        label_parts.append(m['trend'])
                    lp.text = ' '.join(label_parts)
                    lp.font.size = Pt(11)
                    lp.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
                    y_cursor += h + Inches(0.05)
                y_cursor += Inches(0.1)

            elif el_type == 'chart':
                chart_type_str = el.get('chartType', 'bar')
                data_key = el.get('dataKey', '')
                chart_h = Inches(3.5) if size in ('large', 'full') else Inches(2.5)

                # Try to find chart data from the data payload
                charts = data.get('charts', [])
                chart_data_obj = None
                for c in charts:
                    if c.get('dataKey') == data_key or c.get('chartId') == data_key:
                        chart_data_obj = c
                        break
                if not chart_data_obj and charts:
                    chart_data_obj = charts[0]

                if chart_data_obj:
                    cd = CategoryChartData()
                    cd.categories = chart_data_obj.get('categories', [])
                    for series in chart_data_obj.get('series', []):
                        cd.add_series(series.get('name', ''), series.get('data', []))

                    if chart_type_str == 'line':
                        ct = XL_CHART_TYPE.LINE
                    elif chart_type_str == 'pie':
                        ct = XL_CHART_TYPE.PIE
                    else:
                        ct = XL_CHART_TYPE.COLUMN_CLUSTERED
                    if chart_box:
                        chart_x, chart_y, chart_w, chart_h = chart_box
                        # Keep chart values above the template footer / brand rail.
                        chart_h = min(chart_h, sh - chart_y - Inches(1.15))
                    else:
                        chart_x, chart_y, chart_w = margin, y_cursor, content_w
                    chart_frame = slide.shapes.add_chart(ct, chart_x, chart_y, chart_w, chart_h, cd)
                    chart = chart_frame.chart
                    # Slide title already supplies the claim; a duplicate chart
                    # title creates cramped pages and can use a non-CJK fallback.
                    chart.has_title = False
                    chart.has_legend = len(chart_data_obj.get('series', [])) > 1
                    if ct != XL_CHART_TYPE.PIE:
                        chart.value_axis.has_major_gridlines = True
                    for series in chart.series:
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = RGBColor(0xD7, 0x19, 0x20)
                    y_cursor = chart_y + chart_h + Inches(0.08)

            elif el_type == 'comparison':
                entities = el.get('entities', [])
                h = Inches(0.35 * len(entities) + 0.2)
                tf = slide.shapes.add_textbox(margin, y_cursor, content_w, h).text_frame
                tf.word_wrap = True
                for j, ent in enumerate(entities):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    marker = '★ ' if ent.get('highlight') else '  '
                    p.text = f'{marker}{ent.get("name", "")}  {ent.get("value", "")}'
                    p.font.size = Pt(11)
                    if ent.get('highlight'):
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
                    else:
                        p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
                y_cursor += h + Inches(0.1)

            elif el_type == 'table':
                headers = el.get('headers', [])
                rows_data = el.get('rows', [])
                if headers and rows_data:
                    row_count = min(len(rows_data) + 1, 15)  # cap at 15 rows
                    col_count = len(headers)
                    h = Inches(0.35 * row_count)
                    tbl = slide.shapes.add_table(row_count, col_count, margin, y_cursor, content_w, h).table

                    for ci, hdr in enumerate(headers):
                        cell = tbl.cell(0, ci)
                        cell.text = str(hdr)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xC0, 0x39, 0x2B)
                        for p in cell.text_frame.paragraphs:
                            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                            p.font.size = Pt(10)
                            p.font.bold = True

                    for ri, row in enumerate(rows_data[:row_count - 1]):
                        for ci, val in enumerate(row[:col_count]):
                            cell = tbl.cell(ri + 1, ci)
                            cell.text = str(val)
                            for p in cell.text_frame.paragraphs:
                                p.font.size = Pt(10)

                    y_cursor += h + Inches(0.15)

            elif el_type == 'insight':
                h = Inches(0.6)
                tf = slide.shapes.add_textbox(margin, y_cursor, content_w, h).text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f'💡 {content}'
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)
                y_cursor += h + Inches(0.1)

            elif el_type == 'source':
                source_y = min(y_cursor, sh - Inches(0.55)) if chart_element else sh - Inches(0.5)
                tf = slide.shapes.add_textbox(margin, source_y, content_w - Inches(1), Inches(0.25)).text_frame
                p = tf.paragraphs[0]
                p.text = f'資料來源：{content}'
                p.font.size = Pt(7)
                p.font.italic = True
                p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        # Page number (except cover)
        if layout_type != 'cover':
            page_num = slide_spec.get('page', '')
            if page_num and not self._write_placeholder(slide, str(page_num), 'slide_number'):
                tf = slide.shapes.add_textbox(sw - Inches(1), sh - Inches(0.5), Inches(0.6), Inches(0.3)).text_frame
                p = tf.paragraphs[0]
                p.text = str(page_num)
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(0x95, 0xA5, 0xA6)
                p.alignment = PP_ALIGN.RIGHT

    @staticmethod
    def _select_layout(layouts, layout_type):
        def placeholder_types(layout):
            return {placeholder.placeholder_format.type for placeholder in layout.placeholders}

        if layout_type == 'content':
            # Prefer an explicit content layout, then a title+body layout.  A
            # generic "title" match is not enough: title-only pages leave no
            # safe frame for a chart or table.
            for layout in layouts:
                name = (layout.name or '').lower()
                if '內容' in name or 'content' in name:
                    return layout
            for layout in layouts:
                kinds = placeholder_types(layout)
                if PP_PLACEHOLDER.TITLE in kinds and PP_PLACEHOLDER.BODY in kinds:
                    return layout
        elif layout_type == 'section_title':
            for layout in layouts:
                if '章節' in (layout.name or '') or 'section' in (layout.name or '').lower():
                    return layout
        elif layout_type == 'cover':
            for layout in layouts:
                kinds = placeholder_types(layout)
                if PP_PLACEHOLDER.TITLE in kinds and PP_PLACEHOLDER.BODY in kinds:
                    return layout
        elif layout_type == 'backcover':
            for layout in reversed(layouts):
                kinds = placeholder_types(layout)
                if PP_PLACEHOLDER.CENTER_TITLE in kinds or PP_PLACEHOLDER.TITLE in kinds:
                    return layout
        return layouts[0]

    @staticmethod
    def _find_placeholder(slide, types):
        for shape in slide.placeholders:
            if shape.placeholder_format.type in types:
                return shape
        return None

    @staticmethod
    def _remove_shape(shape):
        shape._element.getparent().remove(shape._element)

    @classmethod
    def _write_placeholder(cls, slide, text, role):
        title = slide.shapes.title
        if role == 'title' and title is not None and title.has_text_frame:
            title.text = text
            return True
        if role == 'body':
            shape = cls._find_placeholder(slide, (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE))
            if shape is not None and shape.has_text_frame:
                shape.text = text
                return True
        if role == 'slide_number':
            shape = cls._find_placeholder(slide, (PP_PLACEHOLDER.SLIDE_NUMBER,))
            if shape is not None and shape.has_text_frame:
                shape.text = text
                return True
        return False

    def _cors(self):
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')

    def _error(self, code, msg):
        self.send_response(code)
        self._cors()
        self.send_header('Content-Type', 'application/json')
        self.end_headers()
        self.wfile.write(json.dumps({'error': msg}).encode())
